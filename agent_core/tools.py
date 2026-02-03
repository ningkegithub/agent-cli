import os
import subprocess
import sys
import csv
import io
import textwrap
import base64
from langchain_core.tools import tool
from .utils import INTERNAL_SKILLS_DIR, USER_SKILLS_DIR, get_available_skills_hint, get_skill_suggestions, MEMORY_FILE, ensure_memory_exists, PROJECT_ROOT

# å°è¯•å¯¼å…¥å¯é€‰ä¾èµ–
try:
    import docx
    import pypdf
    import openpyxl
    from pptx import Presentation
    HAS_OFFICE_DEPS = True
except ImportError:
    HAS_OFFICE_DEPS = False

@tool
def run_shell(command: str):
    """æ‰§è¡Œ Shell å‘½ä»¤ã€‚"""
    cmd_stripped = command.strip()
    if cmd_stripped.startswith("python3 ") or cmd_stripped.startswith("python "):
        parts = cmd_stripped.split(" ", 1)
        if len(parts) > 1:
            command = f"{sys.executable} {parts[1]}"
    try:
        result = subprocess.run(command, shell=True, capture_output=True, text=True, timeout=60)
        output = result.stdout
        if result.stderr:
            output += f"\nSTDERR: {result.stderr}"
        if len(output) > 2000:
            output = output[:2000] + "...(truncated)"
        return output
    except Exception as e:
        return f"å‘½ä»¤æ‰§è¡Œé”™è¯¯: {e}"

@tool
def manage_skill(skill_name: str, action: str = "activate"):
    """
    ç®¡ç†ç‰¹æ®ŠæŠ€èƒ½çš„ç”Ÿå‘½å‘¨æœŸã€‚
    
    Args:
        skill_name (str): æŠ€èƒ½åç§°ï¼ˆå¦‚ 'excel_master', 'web_scraper'ï¼‰ã€‚
        action (str): æ“ä½œç±»å‹ã€‚
            - 'activate': (é»˜è®¤) æ¿€æ´»æŠ€èƒ½ã€‚åŠ è½½å…¶ System Prompt æŒ‡ä»¤ï¼Œä½¿å…¶å…·å¤‡ç‰¹å®šé¢†åŸŸçš„ä¸“ä¸šèƒ½åŠ›ã€‚
            - 'deactivate': å¸è½½æŠ€èƒ½ã€‚é‡Šæ”¾ä¸Šä¸‹æ–‡ç©ºé—´ï¼Œé¿å…æŒ‡ä»¤å†²çªæˆ– Token æµªè´¹ã€‚
    """
    action = action.lower()
    normalized_name = skill_name.strip()
    
    # === Action: Activate ===
    if action == "activate":
        search_paths = [
            os.path.join(INTERNAL_SKILLS_DIR, normalized_name, "SKILL.md"),
            os.path.join(USER_SKILLS_DIR, normalized_name, "SKILL.md")
        ]
        target_file = None
        skill_base_dir = None
        
        for path in search_paths:
            if os.path.exists(path):
                target_file = path
                skill_base_dir = os.path.dirname(path)
                break
        
        if target_file and skill_base_dir:
            try:
                with open(target_file, "r", encoding="utf-8") as f:
                    content = f.read()
                injected_content = content.replace("{SKILL_DIR}", skill_base_dir)
                return f"SYSTEM_INJECTION: {injected_content}"
            except Exception as e:
                return f"è¯»å–æŠ€èƒ½æ–‡ä»¶é”™è¯¯: {e}"
        else:
            suggestions = get_skill_suggestions(normalized_name)
            hint = get_available_skills_hint()
            err_msg = f"é”™è¯¯: æœ¬åœ°æœªæ‰¾åˆ°æŠ€èƒ½ '{skill_name}'."
            if suggestions: err_msg += f"ä½ æ˜¯ä¸æ˜¯è¦æ‰¾: {', '.join(suggestions)}."
            if hint: err_msg += f"å¯ç”¨æŠ€èƒ½: {hint}"
            return err_msg

    # === Action: Deactivate ===
    elif action == "deactivate":
        return f"SKILL_DEACTIVATION: {normalized_name}"
    
    else:
        return f"é”™è¯¯: ä¸æ”¯æŒçš„æ“ä½œç±»å‹ '{action}'ã€‚è¯·ä½¿ç”¨ 'activate' æˆ– 'deactivate'ã€‚"

def _read_docx(file_path, outline_only=False):
    doc = docx.Document(file_path)
    lines_pool = []
    outline = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines_pool.append("")
            continue
        if para.style.name.startswith('Heading'):
            try:
                level = int(para.style.name.split()[-1])
                outline.append(f"Line {len(lines_pool) + 1}: {'  ' * (level - 1)}- {text}")
            except: pass
        if len(text) > 120:
            wrapped = textwrap.fill(text, width=120)
            lines_pool.extend(wrapped.splitlines())
        else:
            lines_pool.append(text)
    img_count = len(doc.inline_shapes)
    if img_count > 0: lines_pool.append(f"\n[ç³»ç»Ÿæç¤º: è¯¥æ–‡æ¡£åŒ…å« {img_count} å¼ å›¾ç‰‡]")
    if outline_only:
        if not outline: return "--- æ–‡æ¡£å¤§çº² ---\n[æœªæ£€æµ‹åˆ°æ ‡å‡†æ ‡é¢˜æ ·å¼]\n"
        return "--- æ–‡æ¡£å¤§çº² (ç»“æ„åŒ–å¯¼èˆª) ---\n" + "\n".join(outline)
    return "\n".join(lines_pool)

def _read_pdf(file_path):
    reader = pypdf.PdfReader(file_path)
    text = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        wrapped_lines = []
        for line in page_text.splitlines():
             wrapped_lines.append(textwrap.fill(line, width=120) if len(line) > 120 else line)
        img_count = len(page.images)
        img_placeholder = f"\n[IMAGE_PLACEHOLDER: Page {i+1} åŒ…å« {img_count} å¼ å›¾ç‰‡]\n" if img_count > 0 else ""
        text.append(f"--- Page {i+1} ---\n" + "\n".join(wrapped_lines) + img_placeholder)
    return "\n".join(text)

def _read_excel(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    output = []
    for sheet_name in wb.sheetnames:
        output.append(f"--- Sheet: {sheet_name} ---")
        ws = wb[sheet_name]
        si = io.StringIO()
        writer = csv.writer(si)
        for row in ws.rows:
            writer.writerow([cell.value for cell in row])
        output.append(si.getvalue())
    return "\n".join(output)

def _read_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes: slide_text.append(f"[å¤‡æ³¨]: {notes}")
        if slide_text:
            text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
    return "\n".join(text_content)

@tool
def read_file(file_path: str, start_line: int = 1, end_line: int = -1, outline_only: bool = False):
    """
    å…¨èƒ½æ–‡ä»¶è¯»å–å™¨ã€‚æ”¯æŒè¯»å–çº¯æ–‡æœ¬ (.txt, .md, .py, .json) ä»¥åŠåŠå…¬æ–‡æ¡£ (.docx, .pdf, .xlsx, .pptx)ã€‚
    1. è‡ªåŠ¨è§£æï¼šå¯¹äº Office/PDF/PPT æ–‡æ¡£ï¼Œå·¥å…·ä¼šè‡ªåŠ¨æå–æ–‡æœ¬å†…å®¹åŠå›¾ç‰‡å ä½ç¬¦ä¿¡æ¯ã€‚
    2. å¤§çº²æ¨¡å¼ï¼šè®¾ç½® outline_only=True è·å– Docx ç›®å½•åŠè¡Œå·æ˜ å°„ï¼Œå®ç°ç²¾å‡†è·³è½¬ã€‚
    3. åˆ†é¡µåŠŸèƒ½ï¼šæ”¯æŒä½¿ç”¨ start_line å’Œ end_line è¿›è¡ŒæŒ‰è¡Œåˆ†é¡µè¯»å–ï¼ˆè¡Œå·ä» 1 å¼€å§‹ï¼‰ã€‚
    """
    if not os.path.exists(file_path): return f"é”™è¯¯: æœªæ‰¾åˆ°æ–‡ä»¶ '{file_path}'."
    ext = os.path.splitext(file_path)[1].lower()
    try:
        content_full = ""
        if HAS_OFFICE_DEPS:
            if ext == ".docx":
                content_full = _read_docx(file_path, outline_only=outline_only)
                if outline_only: return content_full
            elif ext == ".pdf": content_full = _read_pdf(file_path)
            elif ext in [".xlsx", ".xls"]: content_full = _read_excel(file_path)
            elif ext == ".pptx": content_full = _read_pptx(file_path)
        
        if not content_full:
            with open(file_path, 'r', encoding='utf-8') as f: lines = f.readlines()
        else:
            lines = content_full.splitlines(keepends=True)
            
        total_lines = len(lines)
        start_idx = max(0, start_line - 1)
        end_idx = min(end_line, total_lines) if end_line != -1 else min(start_idx + 500, total_lines)
        
        selected_lines = lines[start_idx:end_idx]
        header = f"--- æ–‡ä»¶å…ƒæ•°æ® ---\nè·¯å¾„: {file_path}\nè¡Œæ•°: {total_lines} | å½“å‰èŒƒå›´: {start_idx+1}-{end_idx}\n"
        footer = f"\n\n[SYSTEM WARNING]: æ–‡ä»¶æœªè¯»å®Œï¼Œåæ–‡è¿˜æœ‰ {total_lines - end_idx} è¡Œã€‚è¯·è°ƒç”¨ read_file(..., start_line={end_idx+1})." if end_idx < total_lines else ""
        return header + "\n" + "".join(selected_lines) + footer
    except Exception as e: return f"è¯»å–æ–‡ä»¶å‡ºé”™: {e}"

@tool
def write_file(file_path: str, content: str):
    """
    å°†æ–‡æœ¬å†…å®¹å†™å…¥æŒ‡å®šæ–‡ä»¶ï¼ˆå®Œå…¨è¦†ç›–ï¼‰ã€‚
    æ³¨æ„ï¼šç³»ç»Ÿä¼šè‡ªåŠ¨æ‰§è¡Œåˆ†åŒºéš”ç¦»ï¼Œç¡®ä¿æ–‡ä»¶å­˜å‚¨åœ¨æ²™ç®±ç›®å½•ä¸‹ï¼š
    - ä¸šåŠ¡æ•°æ®/æŠ¥è¡¨ (.json, .xlsx, .pptx, .docx, .csv, .txt, .md) -> å¼ºåˆ¶å­˜å…¥ output/
    - ä¸´æ—¶è„šæœ¬/ä»£ç  (.py, .sh) -> å¼ºåˆ¶å­˜å…¥ tmp/
    ä¸¥ç¦ç›´æ¥ä¿®æ”¹é¡¹ç›®æ ¹ç›®å½•æˆ–æ ¸å¿ƒä»£ç åŒºã€‚
    """
    try:
        # 1. è·¯å¾„æ ‡å‡†åŒ–ä¸æå–
        p = file_path.replace("\\", "/")
        name = os.path.basename(p)
        ext = os.path.splitext(name)[1].lower()
        original_path = p
        
        # 2. ä¸¥æ ¼åˆ†åŒºé‡å®šå‘ (Sandboxing)
        if ext in {".py", ".sh", ".js"}:
            # è„šæœ¬æ–‡ä»¶å¼ºåˆ¶é‡å®šå‘åˆ° tmp/
            p = f"tmp/{name}"
        else:
            # ä¸šåŠ¡æ•°æ®æ–‡ä»¶å¼ºåˆ¶é‡å®šå‘åˆ° output/
            p = f"output/{name}"
        
        # 3. æ‰§è¡Œå†™å…¥
        target_path = p
        parent_dir = os.path.dirname(target_path)
        if parent_dir and not os.path.exists(parent_dir): 
            os.makedirs(parent_dir, exist_ok=True)
            
        with open(target_path, 'w', encoding='utf-8') as f: 
            f.write(content)
        
        msg = f"æˆåŠŸå†™å…¥æ–‡ä»¶: {target_path}"
        if target_path != original_path:
            msg += f" (ğŸ›¡ï¸ å·²è§¦å‘åˆ†åŒºéš”ç¦»ï¼Œè‡ªåŠ¨é‡å®šå‘è‡ª {original_path})"
        return msg
    except Exception as e: return f"å†™å…¥æ–‡ä»¶å‡ºé”™: {e}"

@tool
def replace_in_file(file_path: str, old_string: str, new_string: str):
    """ç²¾ç¡®æ›¿æ¢æ–‡ä»¶ä¸­çš„å­—ç¬¦ä¸²ã€‚old_string å¿…é¡»åœ¨æ–‡ä»¶ä¸­å”¯ä¸€å­˜åœ¨ã€‚"""
    if not os.path.exists(file_path): return f"é”™è¯¯: æœªæ‰¾åˆ°æ–‡ä»¶ '{file_path}'."
    try:
        with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
        if old_string not in content: return "é”™è¯¯: åœ¨æ–‡ä»¶ä¸­æœªæ‰¾åˆ° old_string."
        if content.count(old_string) > 1: return "é”™è¯¯: old_string ä¸å”¯ä¸€."
        new_content = content.replace(old_string, new_string)
        with open(file_path, 'w', encoding='utf-8') as f: f.write(new_content)
        return f"æˆåŠŸåœ¨ {file_path} ä¸­å®Œæˆæ›¿æ¢ã€‚"
    except Exception as e: return f"æ›¿æ¢å‡ºé”™: {e}"

@tool
def save_memory(content: str):
    """
    ä¿å­˜é•¿æœŸè®°å¿†ã€‚ç”¨äºè®°å½•ç”¨æˆ·åå¥½ã€é‡è¦äº‹å®æˆ–é•¿æœŸä»»åŠ¡çŠ¶æ€ã€‚
    ä¼šè‡ªåŠ¨è¿›è¡Œç›¸ä¼¼åº¦æ£€æŸ¥ï¼Œé¿å…é‡å¤è®°å½•ã€‚
    """
    try:
        ensure_memory_exists()
        import datetime
        import difflib
        
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        # æ™ºèƒ½å»é‡æ£€æŸ¥
        for line in lines:
            content_part = line[20:].strip() if len(line) > 20 else line.strip()
            if difflib.SequenceMatcher(None, content_part, content.strip()).ratio() > 0.85:
                return f"è®°å¿†å·²å­˜åœ¨ (ç›¸ä¼¼åº¦é«˜)ï¼Œè·³è¿‡å†™å…¥: {content}"
        
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
        entry = f"\n- [{timestamp}] {content}"
        with open(MEMORY_FILE, "a", encoding="utf-8") as f:
            f.write(entry)
        return f"æˆåŠŸä¿å­˜è®°å¿†: {content}"
    except Exception as e: return f"ä¿å­˜è®°å¿†å¤±è´¥: {e}"

@tool
def forget_memory(content: str):
    """
    é—å¿˜é•¿æœŸè®°å¿†ã€‚ç‰©ç†åˆ é™¤ MEMORY.md ä¸­åŒ…å«æŒ‡å®šå…³é”®è¯çš„å†…å®¹ã€‚
    ç”¨äºä¿®æ­£é”™è¯¯æˆ–åˆ é™¤è¿‡æ—¶çš„ç”¨æˆ·åå¥½ã€‚
    """
    try:
        ensure_memory_exists()
        
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        new_lines = [l for l in lines if content not in l or l.startswith('#') or l.startswith('##')]
        deleted_count = len(lines) - len(new_lines)
        
        if deleted_count == 0:
            return f"æœªæ‰¾åˆ°åŒ…å« '{content}' çš„ç›¸å…³è®°å¿†ã€‚"
            
        with open(MEMORY_FILE, "w", encoding="utf-8") as f:
            f.writelines(new_lines)
        return f"æˆåŠŸé—å¿˜ {deleted_count} æ¡ç›¸å…³è®°å¿†ã€‚"
    except Exception as e: return f"é—å¿˜è®°å¿†å¤±è´¥: {e}"

@tool
def search_file(file_path: str, pattern: str, case_sensitive: bool = False):
    """
    åœ¨æ–‡ä»¶ä¸­æœç´¢æŒ‡å®šå…³é”®è¯æˆ–æ­£åˆ™æ¨¡å¼ã€‚æ”¯æŒ Office/PDF/PPT æ–‡æ¡£ã€‚
    """
    if not os.path.exists(file_path): return f"é”™è¯¯: æœªæ‰¾åˆ°æ–‡ä»¶ '{file_path}'."
    ext = os.path.splitext(file_path)[1].lower()
    try:
        content_full = ""
        if HAS_OFFICE_DEPS:
            if ext == ".docx": content_full = _read_docx(file_path)
            elif ext == ".pdf": content_full = _read_pdf(file_path)
            elif ext in [".xlsx", ".xls"]: content_full = _read_excel(file_path)
            elif ext == ".pptx": content_full = _read_pptx(file_path)
        
        lines = content_full.splitlines() if content_full else open(file_path, 'r', encoding='utf-8').readlines()
        matches = []
        import re
        flags = 0 if case_sensitive else re.IGNORECASE
        for i, line in enumerate(lines):
            if re.search(pattern, line, flags):
                matches.append(f"Line {i+1}: {line.strip()}")
        if not matches: return f"æœªæ‰¾åˆ°åŒ¹é… '{pattern}' çš„å†…å®¹."
        count = len(matches)
        result = f"--- æœç´¢ç»“æœ (å…± {count} å¤„åŒ¹é…) ---\n" + "\n".join(matches[:20])
        if count > 20: result += f"\n... (è¿˜æœ‰ {count-20} å¤„åŒ¹é…å·²éšè—)"
        return result
    except Exception as e: return f"æœç´¢å‡ºé”™: {e}"

@tool
def retrieve_knowledge(query: str, collection: str = "documents"):
    """
    è¯­ä¹‰æ£€ç´¢å·¥å…·ã€‚ä»æœ¬åœ°çŸ¥è¯†åº“æˆ–å¯¹è¯å†å²ä¸­æ£€ç´¢ç›¸å…³ä¿¡æ¯ã€‚
    é€‚ç”¨åœºæ™¯ï¼š
    1. æŸ¥é˜…å·²å…¥åº“çš„æ–‡æ¡£ï¼ˆå¦‚ç™½çš®ä¹¦ã€æŠ€æœ¯æ–¹æ¡ˆï¼‰ã€‚ Collection: "documents"
    2. å›å¿†è¿‡å»çš„å¯¹è¯èƒŒæ™¯ï¼ˆæƒ…æ™¯è®°å¿†ï¼‰ã€‚ Collection: "episodic_memory"
    """
    # åŠ¨æ€å®šä½è„šæœ¬
    script_path = os.path.join(INTERNAL_SKILLS_DIR, "knowledge_base/scripts/query.py")
    if not os.path.exists(script_path):
        return "é”™è¯¯: çŸ¥è¯†åº“æŠ€èƒ½è„šæœ¬æœªæ‰¾åˆ°ã€‚"
        
    cmd = [sys.executable, script_path, query, collection]
    try:
        # æ³¨å…¥ PYTHONPATH ç¡®ä¿è„šæœ¬èƒ½æ‰¾åˆ° agent_core
        env = os.environ.copy()
        env["PYTHONPATH"] = PROJECT_ROOT
        
        res = subprocess.run(cmd, capture_output=True, text=True, env=env)
        if res.returncode != 0:
            return f"æ£€ç´¢å¤±è´¥: {res.stderr}"
        return res.stdout
    except Exception as e:
        return f"æ‰§è¡Œé”™è¯¯: {e}"

@tool
def describe_image(image_path: str, prompt: str = "è¯·è¯¦ç»†æè¿°è¿™å¼ å›¾ç‰‡çš„å†…å®¹ã€‚"):
    """
    æŸ¥çœ‹å¹¶åˆ†ææœ¬åœ°å›¾åƒæ–‡ä»¶ï¼ˆæ”¯æŒ PNG, JPG, WEBPï¼‰ã€‚
    æ³¨æ„ï¼šæ­¤å·¥å…·ä¼šç‹¬ç«‹è°ƒç”¨è§†è§‰æ¨¡å‹ (gpt-4o-mini) è¿›è¡Œåˆ†æï¼Œä¸å ç”¨ä¸»æ¨¡å‹çš„ä¸Šä¸‹æ–‡ã€‚
    """
    if not os.path.exists(image_path):
        return f"é”™è¯¯: æœªæ‰¾åˆ°å›¾ç‰‡æ–‡ä»¶ '{image_path}'."
    
    try:
        # è·å–å›¾ç‰‡åç¼€
        ext = os.path.splitext(image_path)[1].lower().strip('.')
        if ext == 'jpg': ext = 'jpeg'
        if ext not in ['png', 'jpeg', 'webp']:
            return f"é”™è¯¯: ä¸æ”¯æŒçš„å›¾ç‰‡æ ¼å¼ '{ext}'ã€‚è¯·ä½¿ç”¨ png, jpg æˆ– webpã€‚"

        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        
        # --- è§†è§‰æ¨¡å‹ç‹¬ç«‹é€šé“ ---
        # æ”¯æŒä»»æ„å…¼å®¹ OpenAI æ¥å£çš„å¤šæ¨¡æ€æ¨¡å‹ (å¦‚ Claude, Gemini, Qwen-VL ç­‰)
        # ä¼˜å…ˆè¯»å– VISION_LLM_X ç³»åˆ—å˜é‡
        from langchain_openai import ChatOpenAI
        from langchain_core.messages import HumanMessage
        
        vision_api_key = os.environ.get("VISION_LLM_API_KEY") or os.environ.get("OPENAI_API_KEY")
        vision_base_url = os.environ.get("VISION_LLM_BASE_URL")
        vision_model = os.environ.get("VISION_LLM_MODEL_NAME") or "gpt-4o-mini"
        
        if not vision_api_key:
            return "é”™è¯¯: è§†è§‰èƒ½åŠ›æœªé…ç½®ã€‚è¯·è®¾ç½® VISION_LLM_API_KEY ç¯å¢ƒå˜é‡ (æ”¯æŒ gpt-4o, claude-3-5, qwen-vl ç­‰å…¼å®¹ OpenAI æ¥å£çš„æ¨¡å‹)ã€‚"
            
        vision_llm = ChatOpenAI(
            model=vision_model,
            api_key=vision_api_key,
            temperature=0,
            base_url=vision_base_url # å¦‚æœä¸º Noneï¼ŒChatOpenAI ä¼šé»˜è®¤ä½¿ç”¨å®˜æ–¹åœ°å€
        )
        
        message = HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{ext};base64,{encoded_string}"},
                },
            ]
        )
        
        response = vision_llm.invoke([message])
        return f"--- å›¾åƒåˆ†æç»“æœ ({image_path}) ---\n[Vision Model: {vision_model}]\n{response.content}"
    except Exception as e:
        return f"å›¾åƒå¤„ç†å‡ºé”™: {e}"

available_tools = [run_shell, manage_skill, read_file, write_file, replace_in_file, search_file, save_memory, forget_memory, retrieve_knowledge, describe_image]
