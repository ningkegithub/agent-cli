import unittest
import os
import subprocess
import shutil
import sys

# 设定测试用的 Collection，避免污染 'documents'
TEST_COLLECTION = "test_integration_rag"
SKILL_DIR = "skills/knowledge_base/scripts"
TEST_DATA_FILE = "tests/test_data/office_mock/test_rag.pptx"
PYTHON_EXE = "./venv/bin/python3"

class TestSkillKnowledgeBase(unittest.TestCase):
    
    @classmethod
    def setUpClass(cls):
        # 确保目录存在
        os.makedirs(os.path.dirname(TEST_DATA_FILE), exist_ok=True)
        
        # 动态生成测试用的 PPT (复用 pptx 库)
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "RAG Test Document"
            slide.placeholders[1].text = "Nebula Core 价格 is 50000 CNY"
            prs.save(TEST_DATA_FILE)
            print(f"✅ Generated mock PPT: {TEST_DATA_FILE}")
        except ImportError:
            raise RuntimeError("python-pptx not installed in test environment")
            
        # 设置基础环境 (不再注入 PYTHONPATH，验证脚本自举能力)
        cls.env = os.environ.copy()
        
        # [新增] 清理旧的测试表，防止 Already Exists 错误
        try:
            # 临时把项目根目录加入 path 以便在此处调用 DBManager
            sys.path.append(os.getcwd())
            from skills.knowledge_base.scripts.db_manager import DBManager
            db = DBManager.get_instance()
            db.reset_table(TEST_COLLECTION)
            print(f"\n🧹 Cleaned up table: {TEST_COLLECTION}")
        except Exception as e:
            print(f"\n⚠️ Cleanup warning: {e}")

    def run_script(self, script_name, args):
        """辅助函数：运行技能脚本"""
        cmd = [PYTHON_EXE, os.path.join(SKILL_DIR, script_name)] + args
        result = subprocess.run(
            cmd, 
            env=self.env, 
            capture_output=True, 
            text=True
        )
        return result

    def test_lifecycle(self):
        print("\n🧪 Testing RAG Lifecycle (Ingest -> Search -> List -> Delete -> Search)...")
        
        # 1. Ingest
        print("  [1/5] Ingesting...")
        res = self.run_script("ingest.py", [TEST_DATA_FILE, TEST_COLLECTION])
        self.assertEqual(res.returncode, 0, f"Ingest failed: {res.stderr}")
        self.assertIn("Ingested", res.stdout)
        
        # 2. Search (Expect Hit)
        print("  [2/5] Searching...")
        res = self.run_script("query.py", ["Nebula Core 价格", TEST_COLLECTION])
        self.assertEqual(res.returncode, 0)
        self.assertIn("50000", res.stdout) # 确保搜到了价格
        self.assertIn("test_rag.pptx", res.stdout) # 确保来源正确 (文件名匹配即可，因为 ingest 会归档重命名，但 source 字段包含文件名)
        
        # 3. List
        print("  [3/5] Listing...")
        res = self.run_script("manage.py", ["list", "--collection", TEST_COLLECTION])
        self.assertEqual(res.returncode, 0)
        self.assertIn("test_rag.pptx", res.stdout)
        
        # 从 list 输出中提取真实的归档文件名 (因为 ingest 加上了 hash)
        # 输出格式: - /path/to/hash_test_rag.pptx (N 片段)
        import re
        match = re.search(r"- (.*test_rag\.pptx)", res.stdout)
        if not match:
            self.fail("Could not find ingested file in list output")
        real_filename = match.group(1).strip()
        print(f"    -> Real filename in DB: {os.path.basename(real_filename)}")
        
        # 4. Delete
        print("  [4/5] Deleting...")
        # 使用提取到的真实路径/文件名进行删除
        res = self.run_script("manage.py", ["delete", real_filename, "--collection", TEST_COLLECTION])
        self.assertEqual(res.returncode, 0)
        self.assertIn("已成功从知识库", res.stdout)
        
        # 5. Search (Expect Miss)
        print("  [5/5] Re-Searching (Verify Deletion)...")
        res = self.run_script("query.py", ["Nebula Core 价格", TEST_COLLECTION])
        # 注意：如果全删空了，可能会报 "知识库不存在或为空"
        is_empty = "知识库" in res.stdout and "为空" in res.stdout
        is_not_found = "未找到" in res.stdout
        self.assertTrue(is_empty or is_not_found, f"Deletion failed? Output: {res.stdout}")
        
        print("  ✅ RAG Lifecycle Test Passed!")

    def test_schema_migration(self):
        """测试数据库 Schema 自动演进/重建"""
        print("\n🧪 Testing Schema Migration...")
        
        # 1. 手动创建一个旧格式的表 (缺少 'location' 字段)
        sys.path.append(os.getcwd())
        from skills.knowledge_base.scripts.db_manager import DBManager
        db = DBManager.get_instance()
        
        MIGRATION_COLLECTION = "test_migration_rag"
        db.reset_table(MIGRATION_COLLECTION)
        
        # 插入一条旧数据
        old_data = [{"vector": [0.1]*384, "text": "Old data", "source": "old.txt", "type": "doc", "line_range": "1-1"}]
        # 注意：这里我们得绕过 create_table 的自动推断，强制创建一个少字段的表？
        # 或者直接用 create_table，只要 data 里没有 location，schema 就不会有 location
        db.create_table(MIGRATION_COLLECTION, data=old_data)
        print("  [1/3] Created old schema table.")
        
        # 验证旧 Schema
        tbl = db.get_table(MIGRATION_COLLECTION)
        self.assertNotIn("location", tbl.schema.names)
        
        # 2. 运行 ingest.py 插入新数据 (包含 location)
        print("  [2/3] Ingesting new data (should trigger migration)...")
        # ingest.py 会读取真实文件，产生带 location 的 data
        res = self.run_script("ingest.py", [TEST_DATA_FILE, MIGRATION_COLLECTION])
        
        # 验证输出中是否有迁移提示 (因为我们在脚本里 print 了)
        # 注意：subprocess capture_output 可能会捕捉不到实时 print，但最终 stdout 会有
        self.assertIn("Auto-migrating", res.stdout)
        self.assertIn("Ingested", res.stdout)
        
        # 3. 验证新 Schema
        print("  [3/3] Verifying new schema...")
        tbl_new = db.get_table(MIGRATION_COLLECTION)
        self.assertIn("location", tbl_new.schema.names)
        
        # 验证新数据是否在里面
        # 由于我们采取的是 Drop & Create，旧数据没了，新数据在
        count = tbl_new.count_rows()
        self.assertGreater(count, 0)
        
        print("  ✅ Schema Migration Test Passed!")

    @classmethod
    def tearDownClass(cls):
        # 清理测试产生的 Collection
        # 由于我们没有暴露 drop_table 接口到 manage.py，这里只能通过 DBManager 内部清理
        # 或者保留着也行，不影响下次测试（因为是 append，或者我们可以先 drop）
        pass

if __name__ == '__main__':
    unittest.main()
